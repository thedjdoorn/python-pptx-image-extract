import argparse
import cv2
import numpy as np
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

parser = argparse.ArgumentParser(description='Extract all images from a PowerPoint presentation')
parser.add_argument('file', help='Path to the input file')
parser.add_argument('-i', '--invert', action='store_true', help='invert the images\' colors')
parser.add_argument('--name', metavar='name', help='specify a name for the output directory (default: same as the file name)')

args = parser.parse_args()

print(args)
if args.name != None:
	output_dir = args.name
else :
	output_dir = ''.join(args.file.split('.')[:-1])

os.mkdir(os.path.abspath(output_dir))

# https://stackoverflow.com/questions/52491656/extracting-images-from-presentation-file
def iter_picture_shapes(prs):
	for slide in prs.slides:
		for shape in slide.shapes:
			if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
				yield shape

i = 0
for picture in iter_picture_shapes(Presentation(args.file)):
	image = picture.image
	image_bytes = image.blob
	output_image = cv2.imdecode(np.frombuffer(image_bytes, np.uint8), cv2.IMREAD_COLOR)
	if args.invert:
		output_image = cv2.bitwise_not(output_image)
	cv2.imwrite(f'./{output_dir}/{i}.jpg', output_image)
	i = i+1